"""Tests for adapting the Student analytics contract to Teacher records."""

import json
import unittest
from unittest.mock import patch
from io import BytesIO

from services.student_data import teacher_records
from services.assignment_generation import generate_draft
from api.routers.materials import _extract_upload


class StudentDataAdapterTest(unittest.TestCase):
    def test_converts_feed_to_existing_teacher_shapes(self):
        feed = {
            "generated_at": "2026-08-19T12:00:00",
            "students": [{
                "student_id": 7,
                "student_code": "s7",
                "name": "Student Seven",
                "average_score": 55,
                "completion_rate": 50,
                "total_submissions": 1,
                "strong_topics": [],
                "weak_topics": ["RAG citations"],
                "recent_submissions": [{"submission_id": 1}],
                "chat_summary": ["Why does grounding matter?"],
            }],
            "topic_metrics": [{
                "topic": "RAG citations",
                "attempts": 2,
                "incorrect": 1,
                "wrong_rate": 50,
            }],
        }
        students, concepts, generated_at = teacher_records(feed)
        self.assertEqual(students[0].id, 7)
        self.assertEqual(students[0].total_submissions, 1)
        self.assertEqual(json.loads(students[0].weak_topics), ["RAG citations"])
        self.assertIn("RAG citations", students[0].recommended_action)
        self.assertEqual(students[0].chat_summary, ["Why does grounding matter?"])
        self.assertEqual(concepts[0].wrong_rate, 50.0)
        self.assertIsNotNone(generated_at)

    @patch("services.assignment_generation.bedrock_client.invoke_json")
    def test_grounded_assignment_generation_validates_result(self, invoke_json):
        invoke_json.return_value = {
            "title": "Grounding checkpoint",
            "question_text": "Explain why evidence is required.",
            "expected_answer": "Evidence connects a claim to course material.",
            "rubric": ["Mentions evidence", "Explains the connection"],
            "source_titles": ["Grounding notes"],
        }
        result = generate_draft(
            "RAG grounding",
            "balanced",
            ["Explain grounding"],
            [{"title": "Grounding notes", "content": "Claims need evidence."}],
        )
        self.assertEqual(result["title"], "Grounding checkpoint")
        prompt = invoke_json.call_args.args[0]
        self.assertIn("Claims need evidence", prompt)

    def test_powerpoint_material_extraction_keeps_slide_locator(self):
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Grounding"
        slide.placeholders[1].text = "Verify every cited claim."
        buffer = BytesIO()
        presentation.save(buffer)
        text, material_type = _extract_upload("lecture.pptx", buffer.getvalue())
        self.assertEqual(material_type, "slide")
        self.assertIn("[Slide 1]", text)
        self.assertIn("Verify every cited claim", text)


if __name__ == "__main__":
    unittest.main()
