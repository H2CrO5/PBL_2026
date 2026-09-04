"""Tests for Teacher material and lecture management."""

import asyncio
from io import BytesIO
import json
import unittest
from unittest.mock import patch

from fastapi import HTTPException, UploadFile
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from api.routers.materials import (
    _extract_upload,
    create_lecture,
    create_material,
    upload_material,
)
from api.schemas.materials import LectureCreateRequest, MaterialCreateRequest
from db.models import Base, Course, Lecture, Material, Teacher
from services.student_data import StudentDataUnavailable


class LectureManagementTest(unittest.TestCase):
    def setUp(self):
        self.engine = create_engine(
            "sqlite://",
            connect_args={"check_same_thread": False},
            poolclass=StaticPool,
        )
        Base.metadata.create_all(self.engine)
        self.session = sessionmaker(bind=self.engine)()
        self.teacher = Teacher(
            teacher_code="teacher-test",
            name="Test Teacher",
            password_hash="not-used",
        )
        self.session.add(self.teacher)
        self.session.flush()
        self.course = Course(
            external_key="course-test",
            teacher_id=self.teacher.id,
            title="Test Course",
            term="2026",
        )
        self.session.add(self.course)
        self.session.commit()

    def tearDown(self):
        self.session.close()
        Base.metadata.drop_all(self.engine)
        self.engine.dispose()

    def _request(self, **overrides):
        values = {
            "course_id": self.course.id,
            "lecture_number": 1,
            "title": "  Grounded Generation  ",
            "learning_objectives": ["  Explain grounding  ", "Cite evidence"],
        }
        values.update(overrides)
        return LectureCreateRequest(**values)

    def test_teacher_can_create_lecture_for_owned_course(self):
        lecture = create_lecture(self._request(), self.teacher, self.session)

        self.assertEqual(lecture.lecture_number, 1)
        self.assertEqual(lecture.title, "Grounded Generation")
        self.assertEqual(
            lecture.learning_objectives,
            ["Explain grounding", "Cite evidence"],
        )

        stored = self.session.get(Course, self.course.id).lectures[0]
        self.assertEqual(json.loads(stored.learning_objectives), lecture.learning_objectives)

    def test_duplicate_lecture_number_is_rejected(self):
        create_lecture(self._request(), self.teacher, self.session)

        with self.assertRaises(HTTPException) as raised:
            create_lecture(
                self._request(title="Another lecture"),
                self.teacher,
                self.session,
            )

        self.assertEqual(raised.exception.status_code, 409)

    def test_teacher_cannot_create_lecture_for_another_teacher_course(self):
        other_teacher = Teacher(
            teacher_code="teacher-other",
            name="Other Teacher",
            password_hash="not-used",
        )
        self.session.add(other_teacher)
        self.session.commit()

        with self.assertRaises(HTTPException) as raised:
            create_lecture(self._request(), other_teacher, self.session)

        self.assertEqual(raised.exception.status_code, 404)

    def test_blank_learning_objectives_are_rejected(self):
        with self.assertRaises(HTTPException) as raised:
            create_lecture(
                self._request(learning_objectives=[" ", ""]),
                self.teacher,
                self.session,
            )

        self.assertEqual(raised.exception.status_code, 422)

    def test_japanese_cp932_text_upload_is_decoded(self):
        text, material_type = _extract_upload(
            "lecture.txt",
            "日本語の講義資料".encode("cp932"),
        )
        self.assertEqual(text, "日本語の講義資料")
        self.assertEqual(material_type, "note")

    def test_powerpoint_text_is_extracted_with_slide_markers(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = "Ground every claim in evidence"
        raw = BytesIO()
        presentation.save(raw)

        text, material_type = _extract_upload("lecture.pptx", raw.getvalue())

        self.assertEqual(material_type, "slide")
        self.assertIn("[Slide 1]", text)
        self.assertIn("Ground every claim in evidence", text)

    def test_password_protected_pdf_has_actionable_error(self):
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        writer.encrypt("secret")
        raw = BytesIO()
        writer.write(raw)

        with self.assertRaises(HTTPException) as raised:
            _extract_upload("protected.pdf", raw.getvalue())

        self.assertEqual(raised.exception.status_code, 422)
        self.assertIn("Password-protected", raised.exception.detail)

    def test_markdown_upload_is_persisted_before_optional_indexing(self):
        create_lecture(self._request(), self.teacher, self.session)
        lecture = self.session.query(Lecture).one()
        upload = UploadFile(
            filename="grounding.md",
            file=BytesIO(b"# Grounding\nUse course evidence."),
        )

        with (
            patch(
                "api.routers.materials.material_storage.store_original",
                return_value="stored/grounding.md",
            ),
            patch(
                "api.routers.materials.student_data.integration_enabled",
                return_value=False,
            ),
        ):
            response = asyncio.run(upload_material(
                course_id=self.course.id,
                lecture_id=lecture.id,
                title="",
                student_visible=True,
                file=upload,
                teacher=self.teacher,
                db=self.session,
            ))

        self.assertEqual(response.title, "grounding")
        self.assertEqual(response.ingestion_status, "local_only")
        self.assertTrue(response.student_visible)
        stored = self.session.query(Material).one()
        self.assertIn("Use course evidence", stored.content)
        self.assertEqual(stored.source_path, "stored/grounding.md")

    @patch("api.routers.materials.student_data.integration_enabled", return_value=True)
    @patch("api.routers.materials.student_data.sync_material")
    def test_material_is_preserved_when_student_indexing_fails(
        self,
        sync_material,
        _integration_enabled,
    ):
        sync_material.side_effect = StudentDataUnavailable("temporary outage")
        create_lecture(self._request(), self.teacher, self.session)
        lecture = self.session.query(Lecture).one()

        response = create_material(
            MaterialCreateRequest(
                course_id=self.course.id,
                lecture_id=lecture.id,
                title="Published notes",
                material_type="note",
                content="Ground claims in evidence.",
                student_visible=True,
            ),
            self.teacher,
            self.session,
        )

        self.assertEqual(response.ingestion_status, "sync_failed")
        self.assertIn("temporary outage", response.sync_error)
        stored = self.session.query(Material).one()
        self.assertEqual(stored.content, "Ground claims in evidence.")
        self.assertTrue(stored.student_visible)


if __name__ == "__main__":
    unittest.main()
