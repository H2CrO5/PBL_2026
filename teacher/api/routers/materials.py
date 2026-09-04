"""Material management endpoints."""

from io import BytesIO
import json
from pathlib import Path

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile, status
from sqlalchemy.orm import Session as DBSession

from api.dependencies import get_current_teacher
from api.schemas.materials import (
    LectureCreateRequest,
    LectureResponse,
    MaterialAudienceRequest,
    MaterialCreateRequest,
    MaterialResponse,
    MaterialSyncAllResponse,
    MaterialSyncResponse,
)
from db.database import get_db
from db.models import Course, Lecture, Material, Teacher
from services import student_data
from services import material_storage
from config import MAX_MATERIAL_UPLOAD_BYTES

router = APIRouter(prefix="/materials", tags=["materials"])


def _student_safe_content(content: str) -> str:
    """Remove explicitly labeled teacher-only sections from public content."""
    public_lines = []
    skipping_internal = False
    internal_prefixes = (
        "teacher note:",
        "teacher note：",
        "teacher prompt:",
        "teacher prompt：",
        "教員メモ:",
        "教員メモ：",
        "教員向けメモ:",
        "教員向けメモ：",
    )
    for line in content.splitlines():
        stripped = line.strip()
        lowered = stripped.lower()
        if any(lowered.startswith(prefix) for prefix in internal_prefixes):
            skipping_internal = True
            continue
        if skipping_internal and (
            stripped.startswith("#")
            or stripped.lower().startswith("[slide ")
            or stripped.lower().startswith("[page ")
        ):
            skipping_internal = False
        if not skipping_internal:
            public_lines.append(line)
    return "\n".join(public_lines).strip()


def _lecture_response(lecture: Lecture) -> LectureResponse:
    return LectureResponse(
        id=lecture.id,
        lecture_number=lecture.lecture_number,
        title=lecture.title,
        learning_objectives=json.loads(lecture.learning_objectives),
    )


def _material_response(material: Material) -> MaterialResponse:
    lecture = material.lecture
    return MaterialResponse(
        id=material.id,
        external_key=material.external_key,
        course_id=material.course_id,
        lecture_id=material.lecture_id,
        lecture_title=lecture.title,
        title=material.title,
        material_type=material.material_type,
        audience=material.audience,
        ingestion_status=material.ingestion_status,
        sync_error=material.sync_error,
        content_preview=material.content[:220],
        created_at=material.created_at,
    )


@router.get("/lectures", response_model=list[LectureResponse])
def get_lectures(
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    course = db.query(Course).filter(Course.teacher_id == teacher.id).first()
    if not course:
        return []

    lectures = db.query(Lecture).filter(Lecture.course_id == course.id).order_by(Lecture.lecture_number).all()
    return [_lecture_response(lecture) for lecture in lectures]


@router.post("/lectures", response_model=LectureResponse)
def create_lecture(
    req: LectureCreateRequest,
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    course = db.query(Course).filter(
        Course.id == req.course_id,
        Course.teacher_id == teacher.id,
    ).first()
    if course is None:
        raise HTTPException(status_code=404, detail="Course not found")

    duplicate = db.query(Lecture).filter(
        Lecture.course_id == course.id,
        Lecture.lecture_number == req.lecture_number,
    ).first()
    if duplicate is not None:
        raise HTTPException(
            status_code=409,
            detail=f"Lecture {req.lecture_number} already exists",
        )

    objectives = [item.strip() for item in req.learning_objectives if item.strip()]
    title = req.title.strip()
    if not title:
        raise HTTPException(status_code=422, detail="Lecture title is required")
    if not objectives:
        raise HTTPException(status_code=422, detail="At least one learning objective is required")

    lecture = Lecture(
        course_id=course.id,
        lecture_number=req.lecture_number,
        title=title,
        learning_objectives=json.dumps(objectives, ensure_ascii=False),
    )
    db.add(lecture)
    db.commit()
    db.refresh(lecture)
    return _lecture_response(lecture)


@router.get("", response_model=list[MaterialResponse])
def get_materials(
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    course = db.query(Course).filter(Course.teacher_id == teacher.id).first()
    if not course:
        return []

    materials = (
        db.query(Material)
        .join(Lecture, Material.lecture_id == Lecture.id)
        .filter(Material.course_id == course.id)
        .order_by(Lecture.lecture_number, Material.id)
        .all()
    )
    return [_material_response(material) for material in materials]


@router.post("", response_model=MaterialResponse)
def create_material(
    req: MaterialCreateRequest,
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    course = db.query(Course).filter(Course.id == req.course_id, Course.teacher_id == teacher.id).first()
    if not course:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Course not found")

    lecture = db.query(Lecture).filter(Lecture.id == req.lecture_id, Lecture.course_id == course.id).first()
    if not lecture:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Lecture not found")

    material = Material(
        course_id=course.id,
        lecture_id=lecture.id,
        title=req.title,
        material_type=req.material_type,
        audience=req.audience,
        content=req.content,
        ingestion_status="local_only",
    )
    db.add(material)
    db.flush()
    material.external_key = f"{course.external_key}:material-{material.id}"
    db.commit()
    db.refresh(material)
    if student_data.integration_enabled():
        _sync_material_record(material, db)
    return _material_response(material)


def _decode_text(content: bytes) -> str:
    if content.startswith((b"\xff\xfe", b"\xfe\xff")):
        return content.decode("utf-16")
    try:
        return content.decode("utf-8-sig")
    except UnicodeDecodeError:
        try:
            return content.decode("cp932")
        except UnicodeDecodeError as exc:
            raise HTTPException(
                status_code=422,
                detail="Text files must use UTF-8, UTF-16, or Japanese CP932 encoding",
            ) from exc


def _extract_upload(filename: str, content: bytes) -> tuple[str, str]:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md"}:
        return _decode_text(content), "note"
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(content))
        if reader.is_encrypted and reader.decrypt("") == 0:
            raise HTTPException(
                status_code=422,
                detail="Password-protected PDF files are not supported",
            )
        pages = [
            f"[Page {index}]\n{page.extract_text() or ''}"
            for index, page in enumerate(reader.pages, start=1)
        ]
        return "\n\n".join(pages), "book"
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(content))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            text_items = [
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            slides.append(f"[Slide {index}]\n" + "\n".join(text_items))
        return "\n\n".join(slides), "slide"
    raise HTTPException(status_code=415, detail="Supported files: PDF, PPTX, MD, TXT")


@router.post("/upload", response_model=MaterialResponse)
async def upload_material(
    course_id: int = Form(...),
    lecture_id: int = Form(...),
    title: str = Form(""),
    audience: str = Form("teacher"),
    file: UploadFile = File(...),
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    course = db.query(Course).filter(
        Course.id == course_id,
        Course.teacher_id == teacher.id,
    ).first()
    lecture = db.query(Lecture).filter(
        Lecture.id == lecture_id,
        Lecture.course_id == course_id,
    ).first()
    if course is None or lecture is None:
        raise HTTPException(status_code=404, detail="Course or lecture not found")
    if audience not in {"student", "teacher"}:
        raise HTTPException(status_code=422, detail="Invalid material audience")
    raw = await file.read(MAX_MATERIAL_UPLOAD_BYTES + 1)
    if not raw:
        raise HTTPException(status_code=422, detail="Material file is empty")
    if len(raw) > MAX_MATERIAL_UPLOAD_BYTES:
        limit_mb = MAX_MATERIAL_UPLOAD_BYTES // (1024 * 1024)
        raise HTTPException(
            status_code=413,
            detail=f"Material file is too large (maximum {limit_mb} MB)",
        )
    try:
        extracted, material_type = _extract_upload(file.filename or "material.txt", raw)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=422, detail="Material file could not be parsed") from exc
    if not extracted.strip():
        suffix = Path(file.filename or "").suffix.lower()
        if suffix == ".pdf":
            detail = "No readable text was found. Scanned or image-only PDFs require OCR and are not supported."
        elif suffix == ".pptx":
            detail = "No readable text was found in the PowerPoint slides."
        else:
            detail = "No readable text could be extracted from this file."
        raise HTTPException(status_code=422, detail=detail)
    try:
        source_path = material_storage.store_original(
            file.filename or "material",
            raw,
            course.external_key or f"course-{course.id}",
        )
    except Exception as exc:
        raise HTTPException(
            status_code=503,
            detail="Original file storage failed; no material was created",
        ) from exc
    material = Material(
        course_id=course.id,
        lecture_id=lecture.id,
        title=title.strip() or Path(file.filename or "Material").stem,
        material_type=material_type,
        audience=audience,
        source_path=source_path,
        content=extracted,
        ingestion_status="local_only",
    )
    db.add(material)
    db.flush()
    material.external_key = f"{course.external_key}:material-{material.id}"
    # Preserve the extracted material before the slower cross-service indexing
    # call. A Student/Bedrock outage must not lose an otherwise valid upload.
    db.commit()
    db.refresh(material)
    if student_data.integration_enabled():
        _sync_material_record(material, db)
    return _material_response(material)


def _sync_payload(material: Material) -> dict:
    content = (
        _student_safe_content(material.content)
        if material.audience == "student"
        else material.content
    )
    audience = material.audience
    if audience == "student" and not content:
        # A document containing only explicitly marked teacher sections has no
        # student-facing material to publish.
        audience = "teacher"
    return {
        "external_material_id": material.external_key,
        "external_course_id": material.lecture.course.external_key,
        "course_title": material.lecture.course.title,
        "term": material.lecture.course.term,
        "lecture_external_id": (
            f"{material.lecture.course.external_key}:lecture-{material.lecture.id}"
        ),
        "lecture_number": material.lecture.lecture_number,
        "lecture_title": material.lecture.title,
        "title": material.title,
        "material_type": material.material_type,
        "audience": audience,
        "content": content or material.content,
    }


def _sync_material_record(material: Material, db: DBSession) -> dict | None:
    """Synchronize one persisted material while retaining actionable state."""
    material.ingestion_status = "indexing"
    material.sync_error = None
    db.commit()
    try:
        result = student_data.sync_material(_sync_payload(material))
    except student_data.StudentDataUnavailable as exc:
        material.ingestion_status = "sync_failed"
        material.sync_error = str(exc)[:500]
        db.commit()
        return None
    material.ingestion_status = result["ingestion_status"]
    material.sync_error = None
    db.commit()
    db.refresh(material)
    return result


@router.post("/sync-all", response_model=MaterialSyncAllResponse)
def sync_all_materials(
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    if not student_data.integration_enabled():
        raise HTTPException(status_code=503, detail="Student integration is not configured")
    materials = (
        db.query(Material)
        .join(Course, Material.course_id == Course.id)
        .filter(Course.teacher_id == teacher.id)
        .all()
    )
    synced = failed = chunks = 0
    for material in materials:
        result = _sync_material_record(material, db)
        if result:
            chunks += result["chunk_count"]
            synced += 1
        else:
            failed += 1
    return MaterialSyncAllResponse(synced=synced, failed=failed, chunks=chunks)


@router.post("/{material_id}/audience", response_model=MaterialResponse)
def update_material_audience(
    material_id: int,
    req: MaterialAudienceRequest,
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    material = (
        db.query(Material)
        .join(Course, Material.course_id == Course.id)
        .filter(Material.id == material_id, Course.teacher_id == teacher.id)
        .first()
    )
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")

    previous_audience = material.audience
    material.audience = req.audience
    material.ingestion_status = "local_only"
    if student_data.integration_enabled():
        result = _sync_material_record(material, db)
        if result is None:
            sync_error = material.sync_error or "Student visibility update failed"
            material.audience = previous_audience
            material.ingestion_status = "sync_failed"
            material.sync_error = sync_error
            db.commit()
            raise HTTPException(status_code=503, detail=sync_error)
    db.commit()
    db.refresh(material)
    return _material_response(material)


@router.post("/{material_id}/sync", response_model=MaterialSyncResponse)
def sync_material(
    material_id: int,
    teacher: Teacher = Depends(get_current_teacher),
    db: DBSession = Depends(get_db),
):
    material = (
        db.query(Material)
        .join(Course, Material.course_id == Course.id)
        .filter(Material.id == material_id, Course.teacher_id == teacher.id)
        .first()
    )
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")
    if not student_data.integration_enabled():
        raise HTTPException(status_code=503, detail="Student integration is not configured")
    result = _sync_material_record(material, db)
    if result is None:
        raise HTTPException(status_code=503, detail=material.sync_error or "Material sync failed")
    return MaterialSyncResponse(
        id=material.id,
        ingestion_status=material.ingestion_status,
        chunk_count=result["chunk_count"],
    )
